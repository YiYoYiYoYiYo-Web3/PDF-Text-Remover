from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import io
import os

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_ppt_from_pages(pages_data, output_path):
    """
    Create a PowerPoint presentation from page data.
    
    Args:
        pages_data: List of dicts, each containing:
            - 'image': PIL Image object (background)
            - 'layout': Layout data from OCR (text_blocks)
        output_path: Path to save the PPTX file
    """
    prs = Presentation()
    
    # Set slide dimensions based on first image
    if pages_data:
        first_image = pages_data[0]['image']
        width_px, height_px = first_image.size
        
        # Convert pixels to inches (assuming 96 DPI)
        prs.slide_width = Inches(width_px / 96.0)
        prs.slide_height = Inches(height_px / 96.0)
    
    for i, page_data in enumerate(pages_data):
        print(f"Creating slide {i+1}/{len(pages_data)}...")
        
        image = page_data['image']
        layout_data = page_data.get('layout', {})
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add background image
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG')
        img_stream.seek(0)
        
        slide.shapes.add_picture(
            img_stream,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # Add text blocks
        text_blocks = layout_data.get('text_blocks', [])
        for block in text_blocks:
            try:
                text = block.get('text', '')
                bbox = block.get('bbox', {})
                font_info = block.get('font', {})
                
                # Get original image dimensions from layout data
                original_size = layout_data.get('original_size', {})
                original_width_px = original_size.get('width', image.size[0])
                original_height_px = original_size.get('height', image.size[1])
                
                # Get current image dimensions
                current_width_px, current_height_px = image.size
                
                # Get original coordinates and size from OCR
                original_x = bbox.get('x', 0)
                original_y = bbox.get('y', 0)
                original_width = bbox.get('width', 100)
                original_height = bbox.get('height', 30)
                
                # Calculate scaling factors
                scale_x = current_width_px / original_width_px
                scale_y = current_height_px / original_height_px
                
                # Apply scaling to coordinates and dimensions
                scaled_x = original_x * scale_x
                scaled_y = original_y * scale_y
                scaled_width = original_width * scale_x
                scaled_height = original_height * scale_y
                
                print(f"[DEBUG PPT] Original size: {original_width_px}x{original_height_px}")
                print(f"[DEBUG PPT] Current size: {current_width_px}x{current_height_px}")
                print(f"[DEBUG PPT] Original: x={original_x}, y={original_y}, w={original_width}, h={original_height}")
                print(f"[DEBUG PPT] Scaled: x={scaled_x}, y={scaled_y}, w={scaled_width}, h={scaled_height}")
                print(f"[DEBUG PPT] Scale factors: x={scale_x:.2f}, y={scale_y:.2f}")
                
                # Convert scaled pixel coordinates to inches
                x = Inches(scaled_x / 96.0)
                y = Inches(scaled_y / 96.0)
                width = Inches(scaled_width / 96.0)
                height = Inches(scaled_height / 96.0)
                
                # Create text box
                textbox = slide.shapes.add_textbox(x, y, width, height)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                text_frame.clear()
                
                p = text_frame.paragraphs[0]
                p.text = text
                
                # Apply font styling
                font = p.font
                font.name = font_info.get('family', 'Arial')
                # Apply scaling to font size (use scale_y for height)                font.size = Pt(font_info.get('size', 12) * scale_y)
                
                # Set font weight
                if font_info.get('weight', 'normal').lower() == 'bold':
                    font.bold = True
                
                # Set font color
                color_hex = font_info.get('color', '#000000')
                try:
                    rgb = hex_to_rgb(color_hex)
                    font.color.rgb = RGBColor(*rgb)
                except:
                    font.color.rgb = RGBColor(0, 0, 0)  # Default to black
                
            except Exception as e:
                print(f"Warning: Failed to add text block: {e}")
                continue
    
    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint saved to {output_path}")
